"""Generate PowerPoint presentation for Canary MCP Server project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(0, 102, 204)  # Blue
SECONDARY_COLOR = RGBColor(51, 51, 51)  # Dark gray
ACCENT_COLOR = RGBColor(255, 153, 0)  # Orange
SUCCESS_COLOR = RGBColor(34, 139, 34)  # Green
BG_COLOR = RGBColor(245, 245, 245)  # Light gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, item in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]

        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(10)

def add_section_divider(prs, section_title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER

# Slide 1: Title
add_title_slide(prs,
    "Universal Canary MCP Server",
    "BD Hackathon 2025-10")

# Slide 2: Problem Statement
add_content_slide(prs, "The Challenge", [
    "🏭 Complex Tag Hierarchies - Deep folder structures with cryptic naming",
    "❌ Access Barriers - Requires API knowledge and custom integration code",
    "⏱️ Time-Consuming - Engineers spend hours navigating and querying data",
    "🔍 Non-Intuitive Naming - Abbreviations and conventions vary by plant",
    "",
    "Example path:",
    "Views/Secil/Portugal/Cement/Maceira/400 - Clinker Production/",
    "431 - Kiln/Normalised/Energy/P431"
])

# Slide 3: The Solution
add_content_slide(prs, "Our Solution", [
    "💬 Natural Language Queries - Ask questions conversationally",
    "🤖 LLM Integration - Connect Claude Desktop to plant data",
    "🔌 MCP Protocol - Standardized bridge between AI and industrial systems",
    "⚡ Instant Access - Seconds instead of hours",
    "",
    "Example: 'What was the average kiln speed yesterday?'",
    "→ System automatically finds tags, retrieves data, analyzes results"
])

# Slide 4: Section divider - Architecture
add_section_divider(prs, "System Architecture")

# Slide 5: High-Level Architecture
add_content_slide(prs, "System Architecture - 3 Layers", [
    "1️⃣ LLM Client Layer",
    "   • Claude Desktop, Continue IDE, other MCP clients",
    "   • Natural language understanding",
    "",
    "2️⃣ MCP Server Layer (Our Implementation)",
    "   • FastMCP framework (Python 3.12+)",
    "   • 5 core tools: ping, list_namespaces, search_tags,",
    "     get_tag_metadata, read_timeseries",
    "   • Authentication & session management",
    "   • Time expression parsing",
    "",
    "3️⃣ Canary Historian Layer",
    "   • SAF authentication",
    "   • Views Web API (Read API)",
    "   • Time-series database"
])

# Slide 6: Authentication Flow
add_content_slide(prs, "Authentication & Session Management", [
    "🔐 Token-Based Authentication with Canary SAF",
    "",
    "Process:",
    "1. User token (from .env) → Session token request",
    "2. Session token cached with 2-minute expiry",
    "3. Automatic refresh when <30 seconds remaining",
    "4. Retry logic with exponential backoff (3 attempts)",
    "",
    "Benefits:",
    "✓ Seamless - Users never see token management",
    "✓ Efficient - Token reuse reduces API calls",
    "✓ Resilient - Auto-retry on connection failures"
])

# Slide 7: Section divider - Features
add_section_divider(prs, "Key Features")

# Slide 8: Core MCP Tools
add_two_column_slide(prs, "5 Core MCP Tools", [
    "✅ ping()",
    "   Health check and connectivity test",
    "",
    "✅ list_namespaces()",
    "   Browse Canary tag hierarchy",
    "   Returns full folder structure",
    "",
    "✅ search_tags(pattern)",
    "   Find tags by wildcard pattern",
    "   Returns matching tags with metadata"
], [
    "✅ get_tag_metadata(tag_path)",
    "   Retrieve tag properties:",
    "   • Data type, units, ranges",
    "   • Description, update rate",
    "",
    "✅ read_timeseries(tags, start, end)",
    "   Fetch historical data:",
    "   • Natural language times",
    "   • Multiple tags supported",
    "   • Quality flags included",
    "   • Pagination handling"
])

# Slide 9: Natural Language Time Parsing
add_content_slide(prs, "Smart Time Expression Parsing", [
    "🕐 Converts Natural Language → ISO 8601 Timestamps",
    "",
    "Supported Expressions:",
    "• 'yesterday' → 2025-10-30T00:00:00Z",
    "• 'past 24 hours' → current time - 24h",
    "• 'last week' → current time - 7 days",
    "• 'now' → current timestamp",
    "",
    "Benefits:",
    "• Users don't need to know exact timestamp formats",
    "• Timezone-aware conversions",
    "• Validation of time ranges (start < end)"
])

# Slide 10: Section divider - Demo
add_section_divider(prs, "Usage & Demo")

# Slide 11: How to Use
add_content_slide(prs, "How to Use the Canary MCP Server", [
    "📦 Installation (3 steps):",
    "1. Install dependencies: uv sync --all-extras",
    "2. Configure .env with Canary credentials",
    "3. Install Claude Desktop config file",
    "",
    "🚀 Usage:",
    "1. Open Claude Desktop (auto-connects to MCP server)",
    "2. Look for 'canary-historian' server indicator",
    "3. Ask questions naturally!",
    "",
    "Example: 'Show me the kiln temperature for yesterday'",
    "→ Claude automatically calls search_tags() and read_timeseries()"
])

# Slide 12: Demo Scenario 1
add_content_slide(prs, "Demo Scenario #1: Quick Status Check", [
    "User Query:",
    "'What is the latest value for kiln 5 431 shell velocity?'",
    "",
    "System Flow:",
    "1. Claude interprets: kiln 431, shell velocity, latest value",
    "2. Calls search_tags('*431*shell*velocity*')",
    "3. Selects best match from candidates",
    "4. Calls read_timeseries(tag, 'past 24 hours', 'now', 1)",
    "",
    "Response:",
    "'The latest shell velocity for Kiln 5 (431) is 1.2 rpm",
    " as of 2025-10-31 14:30 UTC. Quality: Good.'",
    "",
    "⏱️ Time: <5 seconds"
])

# Slide 13: Demo Scenario 2
add_content_slide(prs, "Demo Scenario #2: Trend Analysis", [
    "User Query:",
    "'Compare kiln 431 power consumption yesterday vs today'",
    "",
    "System Flow:",
    "1. Searches for power tag (finds P431)",
    "2. Retrieves yesterday's data (full day)",
    "3. Retrieves today's data (partial)",
    "4. Calculates statistics and compares",
    "",
    "Response:",
    "Yesterday: Avg 1765 kW (min 1720, max 1808)",
    "Today: Avg 1752 kW (min 1735, max 1792)",
    "Observation: -0.7% lower with less variance",
    "",
    "⏱️ Time: <10 seconds"
])

# Slide 14: Demo Scenario 3
add_content_slide(prs, "Demo Scenario #3: Multi-Parameter Correlation", [
    "User Query:",
    "'Show correlation between kiln speed and temperature (past week)'",
    "",
    "System Flow:",
    "1. Searches for speed tag",
    "2. Searches for temperature tag",
    "3. Retrieves 7 days of data for both",
    "4. Analyzes correlation",
    "",
    "Response:",
    "Speed: 1.8-2.1 rpm (avg 1.95)",
    "Temperature: 850-920°C (avg 885)",
    "Correlation: r=0.82 (strong positive)",
    "Insight: Speed >2.0 rpm → temp +15-20°C",
    "",
    "⏱️ Time: <15 seconds"
])

# Slide 15: Section divider - Technical Details
add_section_divider(prs, "Technical Achievements")

# Slide 16: Technical Stack
add_two_column_slide(prs, "Technology Stack & Metrics", [
    "🛠️ Core Technologies:",
    "• Python 3.12+",
    "• FastMCP 0.1.0+ (MCP SDK)",
    "• httpx (async HTTP client)",
    "• uv (package manager)",
    "• pytest (testing framework)",
    "",
    "📊 Code Quality:",
    "• 73% test coverage",
    "• Type-safe with mypy",
    "• Ruff for linting",
    "• ~1,500 lines production code"
], [
    "✨ Key Features:",
    "• Async/await architecture",
    "• Retry with exponential backoff",
    "• Session token auto-refresh",
    "• Comprehensive error handling",
    "• Environment-based config",
    "",
    "📈 Project Stats:",
    "• 5 MCP tools implemented",
    "• 15+ test cases",
    "• 20+ commits",
    "• Complete documentation"
])

# Slide 17: Architecture Highlights
add_content_slide(prs, "Technical Architecture Highlights", [
    "🎯 Design Principles:",
    "",
    "1. Separation of Concerns",
    "   • Auth module (345 lines) - Token management",
    "   • Server module (631 lines) - Tools & parsing",
    "",
    "2. Async-First Design",
    "   • Non-blocking I/O for performance",
    "   • Concurrent API calls where possible",
    "",
    "3. Resilience & Reliability",
    "   • Automatic retry on transient failures",
    "   • Graceful error handling with context",
    "   • Session state management",
    "",
    "4. Developer Experience",
    "   • Type hints throughout",
    "   • Clear error messages",
    "   • Extensive documentation"
])

# Slide 18: Section divider - Results
add_section_divider(prs, "Results & Impact")

# Slide 19: What We Achieved
add_content_slide(prs, "What We Achieved", [
    "✅ Production-Ready MCP Server",
    "   • Complete authentication & session management",
    "   • 5 fully functional tools",
    "   • Comprehensive error handling",
    "",
    "✅ Natural Language Interface",
    "   • Converts conversational queries → API calls",
    "   • Time expression parsing",
    "   • Tag search with wildcards",
    "",
    "✅ Robust Testing",
    "   • 73% code coverage",
    "   • Unit + integration tests",
    "   • Validated against real Canary system",
    "",
    "✅ Complete Documentation",
    "   • README, testing guide, API docs",
    "   • Setup instructions for Windows/Claude Desktop"
])

# Slide 20: Business Impact
add_content_slide(prs, "Business Value & Impact", [
    "⏱️ Time Savings",
    "   From minutes/hours → seconds for data access",
    "",
    "👥 Accessibility",
    "   Enable non-technical users to query plant data",
    "   No need to learn complex APIs or tag structures",
    "",
    "🔍 Faster Insights",
    "   Quick analysis → faster problem resolution",
    "   Easier correlation and trend detection",
    "",
    "📈 Scalability",
    "   Foundation for AI-powered plant operations",
    "   Extensible architecture for future features",
    "",
    "🌍 Multi-Site Potential",
    "   Can be deployed across all Secil plants"
])

# Slide 21: Section divider - Future
add_section_divider(prs, "Future Roadmap")

# Slide 22: Future Enhancements - Phase 1 & 2
add_two_column_slide(prs, "Future Roadmap - Phases 1 & 2", [
    "📍 Phase 1: Semantic Resolution",
    "• Fuzzy matching with confidence",
    "• Context-aware search",
    "  (plant/area/unit)",
    "• Synonym expansion",
    "  (speed ↔ rpm ↔ velocity)",
    "• Multi-language support",
    "  (English/Portuguese)",
    "• Path normalization",
    "",
    "Example:",
    "'velocidade do forno 5'",
    "→ Translates, expands synonyms",
    "→ Ranks by fuzzy match + context"
], [
    "📊 Phase 2: Advanced Analytics",
    "• Statistical aggregations",
    "  (avg, min, max, stddev)",
    "• Trend detection",
    "• Anomaly identification",
    "• Correlation analysis",
    "• Time-series forecasting",
    "• Quality-filtered data",
    "",
    "Example:",
    "'Detect anomalies in kiln",
    " temperature last month'",
    "→ Statistical analysis",
    "→ Highlight outliers"
])

# Slide 23: Future Enhancements - Phase 3 & 4
add_two_column_slide(prs, "Future Roadmap - Phases 3 & 4", [
    "⚡ Phase 3: Performance",
    "• Caching layer",
    "  (browse/search results)",
    "• Connection pooling",
    "• Smart pagination",
    "• Continuation token support",
    "• Rate limiting",
    "• Backpressure handling",
    "",
    "Target:",
    "• <500ms average response",
    "• Support 100+ concurrent users",
    "• 99.9% uptime"
], [
    "🌐 Phase 4: Multi-Site",
    "• Multi-plant configuration",
    "• Site-specific dictionaries",
    "• Cross-site comparison",
    "• Federated queries",
    "• Site context awareness",
    "",
    "Example:",
    "'Compare Outao and Maceira",
    " kiln efficiency this week'",
    "→ Queries both sites",
    "→ Normalizes metrics",
    "→ Presents comparison"
])

# Slide 24: Next Steps
add_content_slide(prs, "Next Steps", [
    "🎯 Immediate (1-2 weeks):",
    "• Pilot deployment with select engineers",
    "• Gather user feedback and usage patterns",
    "• Document common queries and pain points",
    "",
    "🚀 Short-term (1-3 months):",
    "• Implement semantic resolution (fuzzy matching, synonyms)",
    "• Optimize performance (caching, connection pooling)",
    "• Expand test coverage to 85%+",
    "",
    "🌟 Long-term (3-6 months):",
    "• Multi-site rollout across Secil plants",
    "• Advanced analytics capabilities",
    "• Integration with additional LLM clients",
    "• Explore write operations (with strict permissions)"
])

# Slide 25: Conclusion
add_content_slide(prs, "Key Takeaways", [
    "✨ We built a production-ready MCP server connecting LLMs to industrial data",
    "",
    "🎯 Engineers can now query plant data using natural language",
    "",
    "⚡ Data access time reduced from hours → seconds",
    "",
    "🔧 Robust, tested, documented, and ready for deployment",
    "",
    "🚀 Foundation for AI-powered plant operations at Secil",
    "",
    "🌍 Scalable architecture ready for multi-site expansion",
    "",
    "",
    "The future of industrial data access is conversational."
])

# Slide 26: Thank You
add_title_slide(prs,
    "Thank You!",
    "Questions?")

# Save presentation
output_path = "C:\\Github\\BD\\BD-hackaton-2025-10\\Canary_MCP_Server_Presentation.pptx"
prs.save(output_path)

print(f"[OK] Presentation created successfully!")
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
